from sys import argv
import os.path
from string import punctuation
from collections import *
from docx import Document
from pptx import Presentation
import csv
import re
import docx2txt
import pandas as pd
from tkinter import *
from tkinter import ttk
from tkinter.filedialog import askdirectory, askopenfilename, askopenfilenames, asksaveasfilename
from PIL import Image, ImageTk
from itertools import chain
import webbrowser 

### Globals ###
#Names of the source files
filenames = []
#Name of the keywords file
keywords_filename = ""
#Name of the result file
result_filename = ""
#Sentences extracted from the source files
sentences = []
#Keywords extracted from the keywords file
keywords = []
#Message displayed when the analyze succeeds
success_msg = ""
#All the labels
labels = {}
#Labels corresponding to the names of the imported files
labels_file = {}
#All the buttons
buttons = {}
#True if the analyze succeeds and there is a result file, False otherwise
has_result = False

# Get the filename with the keywords
def get_filename():
        # Tk().withdraw()
        # Filetypes allowed
        filetypes = [("Fichiers Excel", "*.xlsx *.xls")]
        print("Initializing Dialogue... \nPlease select a file.")
        filename = askopenfilename(initialdir=os.getcwd(), title='Sélectionnez un fichier', filetypes=filetypes)
        if filename:
                global keywords_filename
                keywords_filename = filename
                return filename 
        else:
                pass

# Returns a String containing the list of the names of the imported files  
def preview_filenames(filenames):
        filenameslabel_str = ""
        global filenameslabel
        for i in range (len(filenames)):
                filenameslabel_str = filenameslabel_str + "\n" + filenames[i].split("/")[-1]
        filenameslabel.set(filenameslabel_str[1:])
        return filenameslabel

# Returns the list of the names of the imported files 
def get_filenames():
        # Tk().withdraw()
        # Filetypes allowed
        filetypes = [("Tous les fichiers", "*.xlsx *.xls *.docx *.pptx *.csv *.txt *.rtf"),("Fichiers Excel","*.xlsx *.xls"),("Fichiers CSV", "*.csv"), ("Fichiers texte", "*.txt *.rtf"), ("Fichiers Word", "*.docx"), ("Fichiers Powerpoint", "*.pptx")]
        print("Initializing Dialogue... \nPlease select a file.")
        tk_filenames = askopenfilenames(initialdir=os.getcwd(), title='Sélectionnez un ou plusieurs fichiers', filetypes=filetypes)
        if tk_filenames:
                filenames_list = list(tk_filenames)
                global filenames
                filenames = filenames_list
                preview_filenames(filenames)
                return filenames_list 
        else:
                pass


# Get the keywords from the file Excel format
def get_keywords():
        keywords_file = get_filename()
        print ("Keyword file to import and read:", keywords_file)
        print ("\nReading file...\n")
        ext = os.path.splitext(keywords_file)[1][1:]
        if ext:
                if ext == "xls" or ext == "xlsx":
                        # Retrieving the data
                        data = pd.read_excel(keywords_file, header = None, encoding="utf-8")
                        # Cleaning the data
                        df = data.dropna(how="all")
                        df = df.replace(u'\xa0',' ')
                        # From Dataframe to Numpy array
                        data_array = df.get_values()
                        cells = []
                        cleanCells = []
                        # Extracting raw data from the dataframe
                        for value in data_array:
                                cells.append(value)
                        # Filtering the useful data (text from the Excel)
                        for cell in cells:
                                for element in cell:
                                        if isinstance(element, str):
                                                element = element.replace(u'\xa0',' ')
                                                cleanCells.append(element)
                        error_msg.set("")
                        global keywords
                        keywords = cleanCells
                        print ("\nFile read finished!")
                        return cleanCells
                else:
                        print ("Please enter a valid file")
                        error_msg.set("Veuillez entrer un fichier au format Excel")
                        return 0
        else:
                pass

# Retrieve all the sentences from the input files
def get_all_sentences():
        filenames_list = get_filenames()
        all_sentences = []
        # Processing file by file
        for one_filename in filenames_list:
                print ("Source file to import and read:", one_filename)
                print ("\nReading file...\n")
                
                # Get the extension of the file
                ext = os.path.splitext(one_filename)[1][1:]

                # Extracting the sentences of the file 

                # Text file
                if ext == "txt" or ext == "rtf":
                        text_file = open(one_filename, 'r')
                        all_lines = text_file.readlines()
                        text_file.close()
                        for line in all_lines:
                                for sentence in line.split("."):
                                        all_sentences.append(sentence)
                        print ("\nFile read finished!")
                
                # CSV file
                elif ext == "csv":
                        text_file = open(one_filename, 'r')
                        all_lines = text_file.readlines()
                        text_file.close()
                        for value in all_lines:
                                for sentence in value.split("."):
                                        all_sentences.append(sentence)
                        print ("\nFile read finished!")

                # -------- WIP .doc file --------
                # # + Allow .doc filetype
                # elif ext == "doc":
                #         docx_file = asksaveasfilename(defaultextension=".docx", title="Choisissez où vous voulez enregistrer votre fichier .docx")
                #         document = Document()
                #         for line in one_filename.readlines()
                #                 document.add_paragraph(line)
                #         document.save(docx_file)
                #         text_data = docx2txt.process(document)
                #         for sentence in text_data.split("."):
                #                 all_sentences.append(sentence)
                #         print ("\nFile read finished!")

                # Word file
                elif ext == "docx": 
                        # Retrieve all the text
                        text_data = docx2txt.process(one_filename)
                        for sentence in text_data.split("."):
                                all_sentences.append(sentence)
                        print ("\nFile read finished!")

                # Excel file
                elif ext == "xls" or ext == "xlsx":
                        # Retrieving the data
                        data = pd.read_excel(one_filename, header = None, encoding='utf-8')
                        # Cleaning the data
                        df = data.dropna(how="all")
                        df = df.replace(u'\xa0',' ')
                        # From Dataframe to Numpy array
                        data_array = df.get_values()
                        cells = []
                        cleanCells = []                        
                        # Extracting raw data from the dataframe
                        for value in data_array:
                                cells.append(value)
                        # Filtering the useful data (text from the Excel)
                        for cell in cells:
                                for element in cell:
                                        if isinstance(element, str):
                                                element = element.replace(u'\xa0',' ')
                                                cleanCells.append(element)
                        for cell in cleanCells:
                                for sentence in str(cell).split("."):
                                        all_sentences.append(sentence)
                        print ("\nFile read finished!")


                # Powerpoint file       
                elif ext =="pptx":
                        prs = Presentation(one_filename)
                        for slide in prs.slides:
                                for shape in slide.shapes:
                                        if shape.has_text_frame:
                                                for sentence in shape.text.split("."):
                                                        sentence = sentence.replace(u'\xa0',' ')
                                                        all_sentences.append(sentence)
                        print ("\nFile read finished!")
        print ("All sentences retrieved!")
        # Optimization : Splitting sentences with no point
        all_sentences = list(flatmap(lambda x: x.splitlines(), all_sentences))
        global sentences
        sentences = all_sentences
        return all_sentences


# Setting of the regex pattern : spaces and/or punctuation between and after the word
def word_matches(word, sentence):
        sentence = "." + sentence + "." 
        pattern = re.compile(r'.*(\s|\W)+' + re.escape(word) + r'(\s|\W)+.*')
        return re.match(pattern, sentence)


# Create the output document and process the matching
def keywords_matching(keywords_list, sentences_list, label_error):
        global filenames, keywords_filename, labels, error_msg
        if keywords_filename == "" or len(filenames) == 0:
                error_msg.set("Veuillez choisir au moins deux fichiers")
                return label_error.update()
        else:
                error_msg.set("")
                label_error.update()
                # Adding the matching sentences to a file
                output_file = asksaveasfilename(defaultextension=".docx", title="Choisissez où vous voulez enregistrer votre fichier de réponses")
                document = Document()
                document.add_heading('Résultats', 0)
                paragraphs = {}
                for idx, word in enumerate(keywords_list):
                        cpt = 0
                        idx = str(idx)
                        paragraphs["key" + idx] = document.add_paragraph('')
                        paragraphs["key" + idx].add_run( "- %s : \n\n" % word).bold = True
                        paragraphs["p" + idx] = document.add_paragraph('')
                        # Case insensitive matching
                        lower_word = word.lower()
                        for sentence in sentences_list: 
                                if word_matches(lower_word, sentence.lower()):
                                        paragraphs["p" + idx].add_run( "\t--> %s.\n\n" % sentence)
                                        cpt +=1
                        paragraphs["key" + idx].add_run(str(cpt) + " occurrences").italic = True
                        document.add_page_break()
                global success_msg, result_filename, has_result
                has_result = True
                result_filename = output_file
                success_msg = 'Fichier de résultats ' + os.path.basename(result_filename) + ' téléchargé !'
                document.save(output_file)
                print (success_msg)
                return output_file

# -------- WIP .doc optimization --------
# # Convert .doc file to .docx
# def convert_to_docx(filename):
#         doc_file = filename
#         docx_file = filename + 'x'
#         if not os.path.exists(docx_file):
#                 os.system('mv ' + doc_file + ' ' + docx_file)
#         else:
#           # already a file with same name as doc exists having docx extension, 
#           # which means it is a different file, so we cant read it
#                 print('Info : file with same name of doc exists having docx extension, so we cant read it')
#         return docx_file

# Homemade flatmap function
def flatmap(f, items):
        return chain.from_iterable(list(map(f, items)))

# Set text to a label
def set_text(label, text):
        label["text"] = text
        return label.update()

# Set multi text to labels --> Set text[1] to label[1], etc.
def set_text_multi(labels, text):
        for i in range(len(filenames)): 
                labels["label_file" + str(i)]["text"] = filenames[i].split("/")[-1]
                labels["label_file" + str(i)].update
        return labels

# Delete all the data created by the user
def reset(mw):
        global sentences, filenames, keywords_filename, keywords, error_msg, success_msg, labels_file, labels, filenameslabel
        # Reset global variables
        filenames = []
        keywords_filename = ""
        result_filename = ""
        sentences = []
        keywords = []
        success_msg = ""
        error_msg.set("")
        filenameslabel.set("")
        has_result = False
        buttons["open_result_btn"].grid_forget()
        # Reset keyword file imported
        labels["label_keyword_file"]["text"] = keywords_filename
        # Reset UI messages
        labels["label_error"]["text"] = error_msg
        labels["label_success"]["text"] = success_msg
        # Reset files imported
        reset_labels(labels_file)
        print ("All data cleaned!")
        return labels_file

# Delete all labels in the input table
def reset_labels(labels):
        for item in labels.values():
                item.destroy()
        labels.clear()
        return labels

# Create a label for each text in the list
def create_labels(mw, text_list, fg):
        global labels_file
        # Reset files imported
        reset_labels(labels_file)     
        if isinstance(text_list, list):
                for i in range (len(text_list)):
                        id_label = "label_file" +  str(i)
                        new_label = Label(mw, text = text_list[i].split("/")[-1], fg = fg)
                        labels_file[id_label] = new_label
                return labels_file
        else:
                pass

# Create a label with textvariable
def create_label_var(mw, id_label, textvar, fg):
        global labels
        new_label = Label(mw, textvariable = textvar, fg = fg)
        labels[id_label] = new_label
        return new_label

# Create a label with static text
def create_label_static(mw, id_label, text, fg):
        global labels
        new_label = Label(mw, text = text, fg = fg)
        labels[id_label] = new_label
        return new_label

# Update the labels
def update_labels(labels):
        for label in labels.values():
                label.update()
        return labels

# Create a button
def create_button(mw, id_button, text, command):
        global buttons
        new_button = ttk.Button(mw, text = text, command = command)
        buttons[id_button] = new_button
        return new_button

# Display result message and button
def display_results():
        global has_result, success_msg
        set_text(labels["label_success"], success_msg)
        if has_result: 
                return buttons["open_result_btn"].grid(row=12, column=1, pady=(0, 15))
        else:
                return buttons["open_result_btn"].grid_forget()


# Graphical User Interface
def main():
        # Init window
        mw = Tk()
        mw.title('Bienvenue sur l\'outil de recherche de mots-clé Kanbios !')
        mw.geometry("850x650")
        ttk.Style().configure("TButton", padding=6, relief="flat", background="#ccc")
        # Themes available ('aqua', 'clam', 'alt', 'default', 'classic')
        ttk.Style().theme_use('clam')
        global labels, error_msg, filenameslabel, result_filename, has_result
        error_msg = StringVar()
        filenameslabel = StringVar()

        # Logo Kanbios
        logo_canvas = Canvas(mw, width=371, height=105)
        logo=PhotoImage(file="logo-kanbios-resized.png")
        logo_canvas.create_image(5, 5,image=logo, anchor="nw")
        logo_canvas.grid(row=1, column=1)

        # Step 1
        create_label_static(mw, 'label_title1', '\n1 - Importer vos mots-clés (fichier Excel)', 'black')
        labels["label_title1"].grid(row=2, column=1)
        # Get the name of the keywords file
        create_button(mw, 'keywords_filename_btn', 'Parcourir...', (lambda : get_keywords() and set_text(labels["label_keyword_file"], keywords_filename.split("/")[-1])))
        buttons["keywords_filename_btn"].grid(row=3, column=1)
        # Preview of the file name
        create_label_var(mw, 'label_keyword_file', keywords_filename, 'blue')
        labels["label_keyword_file"]["text"] = keywords_filename
        labels["label_keyword_file"].grid(row=4, column=1, pady=10)

        # Step 2
        create_label_static(mw, 'label_title2', '\n2 - Importer vos fichiers à analyser', 'black')
        labels["label_title2"].grid(row=5, column=1)
        # Get the names of the files to analyse
        create_button(mw, 'get_filenames_btn', 'Parcourir...', lambda : get_all_sentences() and labels["label_importfile"].update())
        buttons["get_filenames_btn"].grid(row=6, column=1)
        # Preview of the file name [If one label containing all file names]
        create_label_var(mw, 'label_importfile', filenameslabel, 'blue')
        labels["label_importfile"].grid(row=7, column=1, pady=(10,0))
        
        # Step 3
        create_label_static(mw, 'label_title3', '\n\n3 - Télécharger votre fichier de résultats', 'black')
        labels["label_title3"].grid(row=8, column=1)
        # Process the matching
        create_button(mw, 'process_btn', 'Télécharger', lambda : keywords_matching(keywords, sentences, labels["label_error"]) and display_results())
        buttons["process_btn"].grid(row=9, column=1)

        # Success message
        create_label_var(mw, 'label_success', success_msg, 'green')
        labels["label_success"].grid(row=11, column=1, pady=(15, 10))

        # Open the result file
        create_button(mw, 'open_result_btn', 'Ouvrir', lambda : webbrowser.open("file://" + result_filename))

        # Error message
        create_label_var(mw, 'label_error', error_msg, 'red')
        labels["label_error"].grid(row=14, column=1)

        # Reset the data
        create_button(mw, 'reset_btn', 'Réinitialiser', lambda : reset(mw) and update_labels(labels))
        buttons["reset_btn"].grid(row=15, column=1, pady=10)
       

        # Close the window
        create_button(mw, 'close_btn', 'Fermer', mw.destroy)
        buttons["close_btn"].grid(row=16, column=1)

        # Configure the borders of the grid
        mw.grid_rowconfigure(0, weight=1)
        mw.grid_rowconfigure(20, weight=1)
        mw.grid_columnconfigure(0, weight=1)
        mw.grid_columnconfigure(3, weight=1)

        mw.mainloop()

# Main
main()

