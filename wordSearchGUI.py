from sys import argv
import os.path
from string import punctuation
from collections import *
from docx import Document
from pptx import Presentation
import csv
import re
import pandas as pd
from tkinter import *
from tkinter import ttk
from tkinter.filedialog import askdirectory, askopenfilename, askopenfilenames, asksaveasfilename

# Globals
filenames = []
keywords_filename = ""
sentences = []
keywords = []
success_msg = ""

def get_filename():
        # Tk().withdraw()
        print("Initializing Dialogue... \nPlease select a file.")
        filename = askopenfilename(initialdir=os.getcwd(), title='Sélectionnez un fichier')
        global keywords_filename
        keywords_filename = filename
        return filename 

def get_filenames():
        # Tk().withdraw()
        print("Initializing Dialogue... \nPlease select a file.")
        tk_filenames = askopenfilenames(initialdir=os.getcwd(), title='Sélectionnez un ou plusieurs fichiers')
        filenames_list = list(tk_filenames)
        global filenames
        filenames = filenames_list
        return filenames_list 


def get_keywords():
        keywords_file = get_filename()
        print(keywords_file)
        ext = os.path.splitext(keywords_file)[1][1:]
        print(ext)
        if ext == "csv":
                text_file = open(keywords_file, 'r')
                all_lines = text_file.readlines()
                text_file.close()
                all_keywords = []
                for value in all_lines:
                        if value:
                                clean_value = re.sub('[^a-zA-Z0-9]+', '', value)
                                all_keywords.append(clean_value)
                global keywords
                keywords = all_keywords
                print ("\nFile read finished!")
                print (all_keywords)
                return all_keywords
        else:
                print ("Please enter a valid file")
                return 0

def get_all_sentences():
        filenames_list = get_filenames()
        # Processing file by file
        all_sentences = []
        for one_filename in filenames_list:
                print ("Text file to import and read:", one_filename)
                print ("\nReading file...\n")
                
                # Get the extension of the file
                ext = os.path.splitext(one_filename)[1][1:]

                # Extracting the sentences of the file 

                # Text file
                if ext == "txt" or ext == "rtf":
                        text_file = open(one_filename, 'r')
                        all_lines = text_file.readlines()
                        text_file.close()
                        for sentence in all_lines.split("."):
                            all_sentences.append(sentence)
                        print ("\nFile read finished!")
                
                # CSV file
                elif ext == "csv":
                        text_file = open(one_filename, 'r')
                        all_lines = text_file.readlines()
                        text_file.close()
                        for value in all_lines:
                                for sentence in value.split("."):
                                        all_sentences.append(sentence)
                        print ("\nFile read finished!")

                # Word file
                elif ext == "doc" or ext == "docx":
                        text_file = Document(one_filename)
                        all_paragraphs = text_file.paragraphs
                        for paragraph in all_paragraphs:
                                for sentence in paragraph.text.split("."):
                                        all_sentences.append(sentence)
                        print ("\nFile read finished!")

                # Excel file
                elif ext == "xls" or ext == "xlsx":
                        data = pd.read_excel(one_filename, encoding='utf-8')
                        df = data.dropna(how="all")
                        df = df.replace(u'\xa0',' ')
                        # From Dataframe to Numpy array
                        data_array = df.get_values()
                        cells = []
                        cleanCells = []
                        # Extracting raw data from the dataframe
                        for value in data_array:
                                cells.append(value)
                        # Filtering the useful data (text from the Excel)
                        for cell in cells:
                                for element in cell:
                                        if isinstance(element, str):
                                                element = element.replace(u'\xa0',' ')
                                                cleanCells.append(element)
                        # --------
                        # For testing purpose
                        # data_file = open("DataExcel.txt", "w", encoding="utf-8")
                        # for cell in cleanCells:
                        #         data_file.write(str(cell) + "\n\n")
                        # data_file.close()
                        # --------
                        for cell in cleanCells:
                                for sentence in str(cell).split("."):
                                        all_sentences.append(sentence)
                        print ("\nFile read finished!")


                # Powerpoint file       
                elif ext == "ppt" or ext =="pptx":
                        prs = Presentation(one_filename)
                        for slide in prs.slides:
                                for shape in slide.shapes:
                                        if shape.has_text_frame:
                                                for sentence in shape.text.split("."):
                                                        sentence = sentence.replace(u'\xa0',' ')
                                                        all_sentences.append(sentence)
                        print ("\nFile read finished!")

        print ("All sentences retrieved!")
        global sentences
        sentences = all_sentences
        return all_sentences


 # Setting of the regex pattern : spaces and/or punctuation between and after the word
def word_matches(word, sentence):
        sentence = "." + sentence + "." 
        pattern = re.compile(r'.*(\s|\W)+' + re.escape(word) + r'(\s|\W)+.*')
        return re.match(pattern, sentence)


def keywords_matching(keywords_list, sentences_list, label_error):
        global filenames, keywords_filename
        if keywords_filename == "" or len(filenames) == 0:
                label_error["text"] = "Veuillez choisir au moins deux fichiers"
                return label_error.update()
        else:
                label_error["text"] = ""
                label_error.update()
                # Adding the matching sentences to a file
                output_file = open(asksaveasfilename(title='Choisissez où vous voulez enregistrer votre fichier de réponses', defaultextension=".txt"), "w")
                for word in keywords_list:
                        output_file.write( "- %s : \n\n" % word)
                        # Case insensitive matching
                        lowerWord = word.lower()
                        for sentence in sentences_list: 
                                if word_matches(lowerWord, sentence.lower()):
                                        output_file.write( "\t--> %s.\n\n" % sentence)
                global success_msg
                success_msg = 'Fichier de résultats ' + os.path.basename(output_file.name) + ' téléchargé !'
                output_file.close()
                print (success_msg)
                return success_msg

# Set text to a label
def set_text(label, text):
        label["text"] = text
        return label.update()

# Not working
def reset(mw):
        global sentences, filenames, keywords_filename, keywords, success_msg
        filenames = []
        keywords_filename = ""
        sentences = []
        keywords = []
        success_msg = ""
        return sentences, filenames, keywords_filename, keywords, success_msg

def create_labels(mw, text, fg):
        if isinstance(text, str):
                new_label = Label(mw, text = text, fg = fg)
                new_label.pack()
                return new_label.update()
        elif isinstance(text, list):
                # d = {}
                labels = []
                for i in range (len(text)):
                        labels.append(Label(mw, text = text[i].split("/")[-1], fg = fg))
                        # print(d["new_label{0}".format(item)] + '\n')
                for label in labels:
                        print (label["text"])
                        label.pack()
                        label.update()
                # print (d + "\n")
                return labels
        else:
                pass




# GUI
def main():
        mw = Tk()
        mw.title('Welcome to the Kanbios Parser')
        mw.geometry("800x600")
        ttk.Style().configure("TButton", padding=6, relief="flat", background="#ccc")

        # Step 1
        label_title1 = Label(mw, text = '\n1 - Importer vos mots-clés')
        label_title1.pack()
        # Get the name of the keywords file
        keywords_filename_btn = ttk.Button(mw, text = 'Parcourir... (.csv)', command = (lambda : get_keywords() and set_text(label_text1, keywords_filename.split("/")[-1])))
        keywords_filename_btn.pack()
        # Preview of the file name
        label_text1 = Label(mw, text = keywords_filename, fg = 'blue')
        label_text1.pack()

        # Step 2
        label_title2 = Label(mw, text = '\n2 - Importer vos fichiers à analyser')
        label_title2.pack()
        # Get the names of the files to analyse
        # get_filenames_btn = ttk.Button(mw, text = 'Parcourir...', command = lambda : get_all_sentences() and set_text(label_text2, filenames))
        get_filenames_btn = ttk.Button(mw, text = 'Parcourir...', command = lambda : get_all_sentences() and create_labels(mw, filenames, 'blue'))
        get_filenames_btn.pack()
        # Preview of the file name
        label_text2 = Label(mw, text = filenames, fg = 'blue')
        label_text2.pack()

        # Step 3
        label_title3 = Label(mw, text = '\n3 - Télécharger votre fichier de correspondance')
        label_title3.pack()
        # Process the matching
        process_btn = ttk.Button(mw, text = 'Télécharger', command = lambda : keywords_matching(keywords, sentences, label_error) and set_text(label_text3, success_msg))
        process_btn.pack()
        # Success message
        label_text3 = Label(mw, text = success_msg, fg = 'green')
        label_text3.pack()

        # Error message
        label_error = Label(mw, text = "", fg = 'red')
        label_error.pack()

        # Reset the data
        reset_btn = ttk.Button(mw, text = 'Réinitialiser', command = lambda : reset(mw))
        reset_btn.pack(side=BOTTOM)

        # Close the window
        close_btn = ttk.Button(mw, text = 'Fermer', command = mw.destroy)
        close_btn.pack(side=BOTTOM)

        mw.mainloop()

# Main
main()

